import os
import PyPDF2
import docx
import pptx
import openpyxl
import ebooklib.epub
from bs4 import BeautifulSoup
import tkinter as tk
from tkinter import ttk
from tkinter import messagebox
import autogpt
from time import time

data_dir = 'data'

def extract_text():
    start_time = time()
    for filename in os.listdir(data_dir):
        if filename.endswith('.pdf'):
            try:
                with open(os.path.join(data_dir, filename), 'rb') as f:
                    reader = PyPDF2.PdfFileReader(f)
                    text = ''
                    for page in reader.pages:
                        text += page.extract_text()
                    with open(os.path.join(data_dir, f"{filename.replace('.pdf', '.txt')}"), 'w', encoding='utf-8') as f:
                        f.write(text)
            except Exception as e:
                print(f"Error extracting text from {filename}: {e}")
        elif filename.endswith('.docx'):
            try:
                doc = docx.Document(os.path.join(data_dir, filename))
                text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
                with open(os.path.join(data_dir, f"{filename.replace('.docx', '.txt')}"), 'w', encoding='utf-8') as f:
                    f.write(text)
            except Exception as e:
                print(f"Error extracting text from {filename}: {e}")
        elif filename.endswith('.pptx'):
            try:
                prs = pptx.Presentation(os.path.join(data_dir, filename))
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text_runs.append(run.text)
                text = '\n'.join(text_runs)
                with open(os.path.join(data_dir, f"{filename.replace('.pptx', '.txt')}"), 'w', encoding='utf-8') as f:
                    f.write(text)
            except Exception as e:
                print(f"Error extracting text from {filename}: {e}")
        elif filename.endswith('.xlsx'):
            try:
                wb = openpyxl.load_workbook(os.path.join(data_dir, filename))
                text = ''
                for sheetname in wb.sheetnames:
                    sheet = wb[sheetname]
                    for row in sheet.iter_rows():
                        for cell in row:
                            text += str(cell.value) + ' '
                        text += '\n'
                with open(os.path.join(data_dir, f"{filename.replace('.xlsx', '.txt')}"), 'w', encoding='utf-8') as f:
                    f.write(text)
            except Exception as e:
                print(f"Error extracting text from {filename}: {e}")
        elif filename.endswith('.epub'):
            try:
                book = ebooklib.epub.read_epub(os.path.join(data_dir, filename))
                text = ''
                for item in book.get_items():
                    if item.get_type() == ebooklib.ITEM_DOCUMENT:
                        soup = BeautifulSoup(item.get_content(), 'html.parser')
                        text += soup.get_text() + '\n'
                with open(os.path.join(data_dir, f"{filename.replace('.epub', '.txt')}"), 'w', encoding='utf-8') as f:
                    f.write(text)
            except Exception as e:
                print(f"Error extracting text from {filename}: {e}")
    elapsed_time = time() - start_time
    messagebox.showinfo("Text Extraction Complete", f"Text extraction complete in {elapsed_time:.2f} seconds!")

def generate_text():
    start_time = time()
    try:
        autogpt.generate(os.path.join(data_dir, 'generated_text.txt'))
    except Exception as e:
        print(f"Error generating text: {e}")
    elapsed_time = time() - start_time
    messagebox.showinfo("Text Generation Complete", f"Text generation complete in {elapsed_time:.2f} seconds!")

def add_urls():
    pass

def scrape_text():
    pass

root = tk.Tk()
root.title("Web Scraper")
root.geometry("500x300")
root.resizable(False, False)

url_label = ttk.Label(root, text="Enter URLs to scrape:")
url_label.pack(side=tk.TOP, pady=(10, 0))
url_entry = tk.Text(root, height=5)
url_entry.pack(side=tk.TOP, padx=10, pady=(0, 10))

add_button = ttk.Button(root, text="Add URLs", command=add_urls)
add_button.pack(side=tk.LEFT, padx=10)

scrape_button = ttk.Button(root, text="Scrape Text", command=scrape_text)
scrape_button.pack(side=tk.LEFT, padx=10)

extract_button = ttk.Button(root, text="Extract Text", command=extract_text)
extract_button.pack(side=tk.LEFT, padx=10)

generate_button = ttk.Button(root, text="Generate Text", command=generate_text)
generate_button.pack(side=tk.LEFT, padx=10)

progress_bar = ttk.Progressbar(root, mode='determinate')
progress_bar.pack(side=tk.BOTTOM, padx=10, pady=(0, 10))

root.mainloop()
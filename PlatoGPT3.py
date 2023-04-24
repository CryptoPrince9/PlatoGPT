import os
import PyPDF2
import docx
import pptx
import openpyxl
import ebooklib.epub
from bs4 import BeautifulSoup
import tkinter as tk
from tkinter import ttk
from tkinter import messagebox
import autogpt
from time import time

data_dir = 'data'

def extract_text():
    start_time = time()
    for filename in os.listdir(data_dir):
        if filename.endswith('.pdf'):
            try:
                with open(os.path.join(data_dir, filename), 'rb') as f:
                    reader = PyPDF2.PdfFileReader(f)
                    text = ''
                    for page in reader.pages:
                        text += page.extract_text()
                    with open(os.path.join(data_dir, f"{filename.replace('.pdf', '.txt')}"), 'w', encoding='utf-8') as f:
                        f.write(text)
            except Exception as e:
                print(f"Error extracting text from {filename}: {e}")
        elif filename.endswith('.docx'):
            try:
                doc = docx.Document(os.path.join(data_dir, filename))
                text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
                with open(os.path.join(data_dir, f"{filename.replace('.docx', '.txt')}"), 'w', encoding='utf-8') as f:
                    f.write(text)
            except Exception as e:
                print(f"Error extracting text from {filename}: {e}")
        elif filename.endswith('.pptx'):
            try:
                prs = pptx.Presentation(os.path.join(data_dir, filename))
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text_runs.append(run.text)
                text = '\n'.join(text_runs)
                with open(os.path.join(data_dir, f"{filename.replace('.pptx', '.txt')}"), 'w', encoding='utf-8') as f:
                    f.write(text)
            except Exception as e:
                print(f"Error extracting text from {filename}: {e}")
        elif filename.endswith('.xlsx'):
            try:
                wb = openpyxl.load_workbook(os.path.join(data_dir, filename))
                text = ''
                for sheetname in wb.sheetnames:
                    sheet = wb[sheetname]
                    for row in sheet.iter_rows():
                        for cell in row:
                            text += str(cell.value) + ' '
                        text += '\n'
                with open(os.path.join(data_dir, f"{filename.replace('.xlsx', '.txt')}"), 'w', encoding='utf-8') as f:
                    f.write(text)
            except Exception as e:
                print(f"Error extracting text from {filename}: {e}")
        elif filename.endswith('.epub'):
            try:
                book = ebooklib.epub.read_epub(os.path.join(data_dir, filename))
                text = ''
                for item in book.get_items():
                    if item.get_type() == ebooklib.ITEM_DOCUMENT:
                        soup = BeautifulSoup(item.get_content(), 'html.parser')
                        text += soup.text + '\n'
                with open(os.path.join(data_dir, f"{filename.replace('.epub', '.txt')}"), 'w', encoding='utf-8') as f:
                    f.write(text)
            except Exception as e:
                print(f"Error extracting text from {filename}: {e}")
    print(f"Text extraction completed in {time()-start_time:.2f} seconds.")

def summarize_text():
    start_time = time()
    for filename in os.listdir(data_dir):
        if filename.endswith('.txt'):
            try:
                with open(os.path.join(data_dir, filename), 'r', encoding='utf-8') as f:
                    text = f.read()
                    summary = autogpt.summarize(text)
                    with open(os.path.join(data_dir, f"{filename.replace('.txt', '_summary.txt')}"), 'w', encoding='utf-8') as f:
                        f.write(summary)
            except Exception as e:
                print(f"Error summarizing text from {filename}: {e}")
    print(f"Text summarization completed in {time()-start_time:.2f} seconds.")
    
def show_message_box():
    messagebox.showinfo('Processing complete', 'Text extraction and summarization completed.')

root = tk.Tk()
root.title('Text Summarizer')
root.geometry('300x150')

extract_button = ttk.Button(root, text='Extract Text', command=extract_text)
extract_button.pack(pady=10, padx=10, fill='x')

summarize_button = ttk.Button(root, text='Summarize Text', command=summarize_text)
summarize_button.pack(pady=10, padx=10, fill='x')

ok_button = ttk.Button(root, text='OK', command=show_message_box)
ok_button.pack(pady=10, padx=10, fill='x')

root.mainloop()